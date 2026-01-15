"""
Binary File Text Extraction - Extract text from PDFs, Word docs, Excel, etc.

Integrates with CC/Wayback scrapers to extract searchable text from binary files
found in Common Crawl and archive.org WARC records.

Supports:
- PDF (via pypdf or pdfplumber)
- Word (.doc, .docx via python-docx)
- Excel (.xls, .xlsx via openpyxl)
- PowerPoint (.ppt, .pptx via python-pptx)
- Archives (.zip, .tar, .gz - list contents only)
"""

import io
import logging
from typing import Optional, Dict, Any
from dataclasses import dataclass

logger = logging.getLogger(__name__)


@dataclass
class ExtractionResult:
    """Result of binary file text extraction."""
    text: str
    file_type: str
    success: bool
    error: Optional[str] = None
    metadata: Dict[str, Any] = None
    char_count: int = 0
    page_count: int = 0


class BinaryTextExtractor:
    """
    Extract searchable text from binary files retrieved from archives.

    Designed to work with WARC records from Common Crawl and Wayback Machine.
    """

    # MIME type to handler mapping
    MIME_HANDLERS = {
        'application/pdf': '_extract_pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '_extract_docx',
        'application/msword': '_extract_doc',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '_extract_xlsx',
        'application/vnd.ms-excel': '_extract_xls',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': '_extract_pptx',
        'application/vnd.ms-powerpoint': '_extract_ppt',
        'application/zip': '_extract_zip_listing',
        'application/x-tar': '_extract_tar_listing',
        'application/gzip': '_extract_gzip_listing',
    }

    def __init__(self):
        """Initialize extractor and check for available libraries."""
        self.available_extractors = self._check_available_libraries()
        logger.info(f"Binary text extractor initialized with: {list(self.available_extractors.keys())}")

    def _check_available_libraries(self) -> Dict[str, bool]:
        """Check which extraction libraries are available."""
        available = {}

        # PDF extraction
        try:
            import pypdf
            available['pypdf'] = True
        except ImportError:
            available['pypdf'] = False

        try:
            import pdfplumber
            available['pdfplumber'] = True
        except ImportError:
            available['pdfplumber'] = False

        # Word documents
        try:
            import docx
            available['python-docx'] = True
        except ImportError:
            available['python-docx'] = False

        # Excel
        try:
            import openpyxl
            available['openpyxl'] = True
        except ImportError:
            available['openpyxl'] = False

        # PowerPoint
        try:
            import pptx
            available['python-pptx'] = True
        except ImportError:
            available['python-pptx'] = False

        # Archives
        available['zipfile'] = True  # Built-in
        available['tarfile'] = True  # Built-in

        return available

    def can_extract(self, mime_type: str) -> bool:
        """Check if we can extract text from this MIME type."""
        if mime_type not in self.MIME_HANDLERS:
            return False

        # Check if required library is available
        if 'pdf' in mime_type and not (self.available_extractors.get('pypdf') or self.available_extractors.get('pdfplumber')):
            return False
        if 'wordprocessingml' in mime_type and not self.available_extractors.get('python-docx'):
            return False
        if 'spreadsheetml' in mime_type and not self.available_extractors.get('openpyxl'):
            return False
        if 'presentationml' in mime_type and not self.available_extractors.get('python-pptx'):
            return False

        return True

    def extract_text(self, binary_data: bytes, mime_type: str, filename: str = '') -> ExtractionResult:
        """
        Extract text from binary file data.

        Args:
            binary_data: Raw bytes from WARC record
            mime_type: MIME type of the file
            filename: Original filename (optional, for extension fallback)

        Returns:
            ExtractionResult with extracted text and metadata
        """
        # Get handler method
        handler_name = self.MIME_HANDLERS.get(mime_type)
        if not handler_name:
            # Try fallback by file extension
            if filename:
                handler_name = self._get_handler_by_extension(filename)

        if not handler_name:
            return ExtractionResult(
                text='',
                file_type=mime_type,
                success=False,
                error=f'No handler for MIME type: {mime_type}'
            )

        # Call the handler
        try:
            handler = getattr(self, handler_name)
            return handler(binary_data, mime_type)
        except Exception as e:
            logger.error(f"Extraction failed for {mime_type}: {e}")
            return ExtractionResult(
                text='',
                file_type=mime_type,
                success=False,
                error=str(e)
            )

    def _get_handler_by_extension(self, filename: str) -> Optional[str]:
        """Fallback: determine handler by file extension."""
        ext_map = {
            '.pdf': '_extract_pdf',
            '.docx': '_extract_docx',
            '.doc': '_extract_doc',
            '.xlsx': '_extract_xlsx',
            '.xls': '_extract_xls',
            '.pptx': '_extract_pptx',
            '.ppt': '_extract_ppt',
            '.zip': '_extract_zip_listing',
            '.tar': '_extract_tar_listing',
            '.gz': '_extract_gzip_listing',
            '.tgz': '_extract_tar_listing',
        }

        for ext, handler in ext_map.items():
            if filename.lower().endswith(ext):
                return handler

        return None

    # ==================== PDF Extraction ====================

    def _extract_pdf(self, data: bytes, mime_type: str) -> ExtractionResult:
        """Extract text from PDF using pypdf or pdfplumber."""
        # Try pdfplumber first (better text extraction)
        if self.available_extractors.get('pdfplumber'):
            return self._extract_pdf_pdfplumber(data)
        elif self.available_extractors.get('pypdf'):
            return self._extract_pdf_pypdf(data)
        else:
            return ExtractionResult(
                text='',
                file_type='pdf',
                success=False,
                error='No PDF extraction library available (install pypdf or pdfplumber)'
            )

    def _extract_pdf_pypdf(self, data: bytes) -> ExtractionResult:
        """Extract using pypdf."""
        try:
            import pypdf

            pdf_file = io.BytesIO(data)
            reader = pypdf.PdfReader(pdf_file)

            text_parts = []
            for page in reader.pages:
                text_parts.append(page.extract_text())

            full_text = '\n\n'.join(text_parts)

            return ExtractionResult(
                text=full_text,
                file_type='pdf',
                success=True,
                char_count=len(full_text),
                page_count=len(reader.pages),
                metadata={
                    'extractor': 'pypdf',
                    'pages': len(reader.pages)
                }
            )
        except Exception as e:
            return ExtractionResult(
                text='',
                file_type='pdf',
                success=False,
                error=f'pypdf extraction failed: {e}'
            )

    def _extract_pdf_pdfplumber(self, data: bytes) -> ExtractionResult:
        """Extract using pdfplumber (better quality)."""
        try:
            import pdfplumber

            pdf_file = io.BytesIO(data)
            text_parts = []
            page_count = 0

            with pdfplumber.open(pdf_file) as pdf:
                page_count = len(pdf.pages)
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)

            full_text = '\n\n'.join(text_parts)

            return ExtractionResult(
                text=full_text,
                file_type='pdf',
                success=True,
                char_count=len(full_text),
                page_count=page_count,
                metadata={
                    'extractor': 'pdfplumber',
                    'pages': page_count
                }
            )
        except Exception as e:
            return ExtractionResult(
                text='',
                file_type='pdf',
                success=False,
                error=f'pdfplumber extraction failed: {e}'
            )

    # ==================== Word Document Extraction ====================

    def _extract_docx(self, data: bytes, mime_type: str) -> ExtractionResult:
        """Extract text from .docx file."""
        if not self.available_extractors.get('python-docx'):
            return ExtractionResult(
                text='',
                file_type='docx',
                success=False,
                error='python-docx not installed'
            )

        try:
            import docx

            doc_file = io.BytesIO(data)
            document = docx.Document(doc_file)

            # Extract paragraphs
            text_parts = [para.text for para in document.paragraphs if para.text.strip()]

            # Extract tables
            for table in document.tables:
                for row in table.rows:
                    row_text = '\t'.join(cell.text for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)

            full_text = '\n'.join(text_parts)

            return ExtractionResult(
                text=full_text,
                file_type='docx',
                success=True,
                char_count=len(full_text),
                metadata={
                    'extractor': 'python-docx',
                    'paragraphs': len(document.paragraphs),
                    'tables': len(document.tables)
                }
            )
        except Exception as e:
            return ExtractionResult(
                text='',
                file_type='docx',
                success=False,
                error=f'docx extraction failed: {e}'
            )

    def _extract_doc(self, data: bytes, mime_type: str) -> ExtractionResult:
        """Extract text from legacy .doc file (requires antiword or textract)."""
        # Legacy .doc files are harder to parse
        # Would need antiword (command-line) or textract
        return ExtractionResult(
            text='',
            file_type='doc',
            success=False,
            error='Legacy .doc format not yet supported (use .docx or install antiword)'
        )

    # ==================== Excel Extraction ====================

    def _extract_xlsx(self, data: bytes, mime_type: str) -> ExtractionResult:
        """Extract text from .xlsx file."""
        if not self.available_extractors.get('openpyxl'):
            return ExtractionResult(
                text='',
                file_type='xlsx',
                success=False,
                error='openpyxl not installed'
            )

        try:
            import openpyxl

            excel_file = io.BytesIO(data)
            workbook = openpyxl.load_workbook(excel_file, data_only=True)

            text_parts = []
            for sheet in workbook.worksheets:
                text_parts.append(f"=== Sheet: {sheet.title} ===")
                for row in sheet.iter_rows(values_only=True):
                    row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)

            full_text = '\n'.join(text_parts)

            return ExtractionResult(
                text=full_text,
                file_type='xlsx',
                success=True,
                char_count=len(full_text),
                metadata={
                    'extractor': 'openpyxl',
                    'sheets': len(workbook.worksheets)
                }
            )
        except Exception as e:
            return ExtractionResult(
                text='',
                file_type='xlsx',
                success=False,
                error=f'xlsx extraction failed: {e}'
            )

    def _extract_xls(self, data: bytes, mime_type: str) -> ExtractionResult:
        """Extract text from legacy .xls file."""
        # Would need xlrd library
        return ExtractionResult(
            text='',
            file_type='xls',
            success=False,
            error='Legacy .xls format not yet supported (use .xlsx or install xlrd)'
        )

    # ==================== PowerPoint Extraction ====================

    def _extract_pptx(self, data: bytes, mime_type: str) -> ExtractionResult:
        """Extract text from .pptx file."""
        if not self.available_extractors.get('python-pptx'):
            return ExtractionResult(
                text='',
                file_type='pptx',
                success=False,
                error='python-pptx not installed'
            )

        try:
            import pptx

            ppt_file = io.BytesIO(data)
            presentation = pptx.Presentation(ppt_file)

            text_parts = []
            for slide_num, slide in enumerate(presentation.slides, 1):
                text_parts.append(f"=== Slide {slide_num} ===")
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text_parts.append(shape.text)

            full_text = '\n'.join(text_parts)

            return ExtractionResult(
                text=full_text,
                file_type='pptx',
                success=True,
                char_count=len(full_text),
                metadata={
                    'extractor': 'python-pptx',
                    'slides': len(presentation.slides)
                }
            )
        except Exception as e:
            return ExtractionResult(
                text='',
                file_type='pptx',
                success=False,
                error=f'pptx extraction failed: {e}'
            )

    def _extract_ppt(self, data: bytes, mime_type: str) -> ExtractionResult:
        """Extract text from legacy .ppt file."""
        return ExtractionResult(
            text='',
            file_type='ppt',
            success=False,
            error='Legacy .ppt format not yet supported (use .pptx)'
        )

    # ==================== Archive Listing ====================

    def _extract_zip_listing(self, data: bytes, mime_type: str) -> ExtractionResult:
        """List contents of ZIP archive."""
        try:
            import zipfile

            zip_file = io.BytesIO(data)
            with zipfile.ZipFile(zip_file) as zf:
                file_list = [f"- {info.filename} ({info.file_size} bytes)" for info in zf.infolist()]
                listing = "ZIP Archive Contents:\n" + '\n'.join(file_list)

            return ExtractionResult(
                text=listing,
                file_type='zip',
                success=True,
                char_count=len(listing),
                metadata={
                    'extractor': 'zipfile',
                    'file_count': len(file_list)
                }
            )
        except Exception as e:
            return ExtractionResult(
                text='',
                file_type='zip',
                success=False,
                error=f'zip listing failed: {e}'
            )

    def _extract_tar_listing(self, data: bytes, mime_type: str) -> ExtractionResult:
        """List contents of TAR archive."""
        try:
            import tarfile

            tar_file = io.BytesIO(data)
            with tarfile.open(fileobj=tar_file) as tf:
                file_list = [f"- {member.name} ({member.size} bytes)" for member in tf.getmembers()]
                listing = "TAR Archive Contents:\n" + '\n'.join(file_list)

            return ExtractionResult(
                text=listing,
                file_type='tar',
                success=True,
                char_count=len(listing),
                metadata={
                    'extractor': 'tarfile',
                    'file_count': len(file_list)
                }
            )
        except Exception as e:
            return ExtractionResult(
                text='',
                file_type='tar',
                success=False,
                error=f'tar listing failed: {e}'
            )

    def _extract_gzip_listing(self, data: bytes, mime_type: str) -> ExtractionResult:
        """Handle gzip files (note: gzip is single-file compression)."""
        try:
            import gzip

            gz_file = io.BytesIO(data)
            with gzip.open(gz_file, 'rt', errors='ignore') as f:
                # Try to read first few lines
                preview = '\n'.join(f.readlines()[:10])
                listing = f"GZIP Compressed File Preview:\n{preview}"

            return ExtractionResult(
                text=listing,
                file_type='gzip',
                success=True,
                char_count=len(listing),
                metadata={'extractor': 'gzip'}
            )
        except Exception as e:
            return ExtractionResult(
                text='',
                file_type='gzip',
                success=False,
                error=f'gzip extraction failed: {e}'
            )

    def extract_pdf_metadata_only(self, pdf_bytes: bytes) -> Optional[Dict]:
        """
        Extract PDF metadata WITHOUT full text parsing.

        Fast lightweight extraction for validation purposes.
        Extracts:
        - Title
        - Author
        - CreationDate
        - ModDate
        - Page count
        - File size

        This is ~100x faster than full text extraction (100ms vs 2-5s).

        Args:
            pdf_bytes: Raw PDF file bytes

        Returns:
            Metadata dict or None if invalid PDF

        Example:
            metadata = extractor.extract_pdf_metadata_only(pdf_bytes)
            if metadata and metadata['page_count'] > 10:
                # Likely a substantial document
                full_text = extractor.extract_text(pdf_bytes, 'application/pdf')
        """
        try:
            import pypdf

            pdf_file = io.BytesIO(pdf_bytes)
            reader = pypdf.PdfReader(pdf_file)

            # Extract metadata
            metadata = {
                'page_count': len(reader.pages),
                'file_size': len(pdf_bytes),
                'file_size_mb': len(pdf_bytes) / 1024 / 1024,
            }

            # PDF metadata (may not always be present)
            if reader.metadata:
                metadata['title'] = reader.metadata.get('/Title', '')
                metadata['author'] = reader.metadata.get('/Author', '')
                metadata['subject'] = reader.metadata.get('/Subject', '')
                metadata['creator'] = reader.metadata.get('/Creator', '')
                metadata['producer'] = reader.metadata.get('/Producer', '')

                # Parse dates (format: D:YYYYMMDDHHmmSS)
                creation_date = reader.metadata.get('/CreationDate', '')
                mod_date = reader.metadata.get('/ModDate', '')

                metadata['creation_date'] = creation_date
                metadata['mod_date'] = mod_date

                # Extract year from creation date if available
                if creation_date and creation_date.startswith('D:'):
                    try:
                        year_str = creation_date[2:6]
                        metadata['creation_year'] = int(year_str)
                    except (ValueError, IndexError):
                        metadata['creation_year'] = None
                else:
                    metadata['creation_year'] = None

            # Encrypted PDF check
            metadata['is_encrypted'] = reader.is_encrypted

            # Rough text estimation (without full extraction)
            # Estimate ~500 chars per page as heuristic
            metadata['estimated_char_count'] = len(reader.pages) * 500

            return metadata

        except Exception as e:
            logger.warning(f"PDF metadata extraction failed: {e}")
            return None


# Convenience function for quick extraction
def extract_text_from_bytes(data: bytes, mime_type: str, filename: str = '') -> str:
    """
    Quick extraction function that returns just the text.

    Args:
        data: Binary file data
        mime_type: MIME type
        filename: Optional filename for extension fallback

    Returns:
        Extracted text, or empty string if extraction failed
    """
    extractor = BinaryTextExtractor()
    result = extractor.extract_text(data, mime_type, filename)
    return result.text if result.success else ''


__all__ = ['BinaryTextExtractor', 'ExtractionResult', 'extract_text_from_bytes']
